# -*- coding: utf-8 -*-
"""Generate MMS PPT — 24 slides, text loaded from JSON"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = os.path.dirname(os.path.abspath(__file__))
OUTPUT = os.path.join(ROOT, "MMS_.pptx")

data = json.load(open(os.path.join(ROOT, "mms_ppt_data.json"), encoding="utf-8"))

# Colors
NAVY = RGBColor(0x1B, 0x3A, 0x5C)
BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE = RGBColor(0xD6, 0xEB, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
LIGHT_GRAY2 = RGBColor(0xE8, 0xE8, 0xE8)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
LIGHT_GREEN = RGBColor(0xD5, 0xF5, 0xE3)
LIGHT_RED = RGBColor(0xFD, 0xED, 0xEC)
LIGHT_ORANGE = RGBColor(0xFE, 0xF3, 0xCB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_box(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    s.line.fill.background()
    if line:
        s.line.fill.solid(); s.line.color.rgb = line; s.line.width = Pt(1.5)
    return s

def add_tb(slide, l, t, w, h):
    return slide.shapes.add_textbox(l, t, w, h)

def txt(tf, text, size=16, bold=False, color=DARK, font="Microsoft YaHei", align=PP_ALIGN.LEFT):
    tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return p

def add_p(tf, text, size=14, bold=False, color=DARK, font="Microsoft YaHei", align=PP_ALIGN.LEFT, sb=Pt(4)):
    p = tf.add_paragraph(); p.alignment = align; p.space_before = sb
    r = p.add_run(); r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return p

def new_slide(title="", dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    if dark:
        add_bg(slide, NAVY)
    else:
        add_bg(slide, WHITE)
        add_box(slide, 0, 0, W, Pt(5), fill=BLUE)
    # number
    n = len(prs.slides)
    tb = add_tb(slide, W - Inches(0.8), H - Inches(0.4), Inches(0.6), Inches(0.3))
    txt(tb.text_frame, str(n), 10, color=GRAY, align=PP_ALIGN.RIGHT)
    if title:
        c = WHITE if dark else DARK
        tb = add_tb(slide, Inches(0.6), Inches(0.25), Inches(12), Inches(0.6))
        txt(tb.text_frame, title, 28, True, c)
        add_box(slide, Inches(0.6), Inches(0.85), Inches(3), Pt(3), fill=GREEN if dark else BLUE)
    return slide

def add_table(slide, l, t, w, h, rows, cols, data, cw=None):
    ts = slide.shapes.add_table(rows, cols, l, t, w, h)
    tbl = ts.table
    if cw:
        for i, cw_ in enumerate(cw):
            tbl.columns[i].width = cw_
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ""
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
            val = str(data[r][c]) if r < len(data) and c < len(data[r]) else ""
            run = p.add_run(); run.text = val
            run.font.name = "Microsoft YaHei"
            if r == 0:
                run.font.size = Pt(11); run.font.bold = True; run.font.color.rgb = WHITE
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            else:
                run.font.size = Pt(10); run.font.color.rgb = DARK
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Pt(6); cell.margin_right = Pt(6)
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3)
    return ts

def code_box(slide, l, t, w, h, text, bg=LIGHT_GRAY):
    add_box(slide, l, t, w, h, fill=bg)
    tb = add_tb(slide, l + Inches(0.15), t + Inches(0.08), w - Inches(0.3), h - Inches(0.16))
    txt(tb.text_frame, text, 10, color=DARK, font="Consolas")

def card(slide, l, t, w, h, title, body, accent=BLUE, bg=WHITE):
    add_box(slide, l, t, w, h, fill=bg, line=RGBColor(0xDD, 0xDD, 0xDD))
    tb = add_tb(slide, l + Inches(0.15), t + Inches(0.08), w - Inches(0.3), h - Inches(0.16))
    txt(tb.text_frame, title, 13, True, accent)
    add_p(tb.text_frame, body, 11, color=DARK)

def card_with_icon(slide, l, t, w, h, icon, title, body, accent=BLUE):
    add_box(slide, l, t, w, h, fill=WHITE, line=RGBColor(0xDD, 0xDD, 0xDD))
    add_box(slide, l, t, Pt(4), h, fill=accent)
    tb = add_tb(slide, l + Inches(0.25), t + Inches(0.08), w - Inches(0.4), h - Inches(0.16))
    txt(tb.text_frame, icon + "  " + title, 13, True, DARK)
    add_p(tb.text_frame, body, 11, color=GRAY)

# ========== SLIDES ==========

# 1 - Cover
slide = new_slide(dark=True)
tb = add_tb(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5))
txt(tb.text_frame, "MMS 测试：工业 CFD 软件验证的\n方法论与实践", 34, True, WHITE, align=PP_ALIGN.CENTER)
tb = add_tb(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6))
txt(tb.text_frame, "Manufactured Solutions \u2014 PDE 求解器的\u201c单元测试\u201d", 18, color=RGBColor(0xBB, 0xDD, 0xFF), align=PP_ALIGN.CENTER)
add_box(slide, Inches(5), Inches(4.5), Inches(3.3), Pt(3), fill=GREEN)
tb = add_tb(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5))
txt(tb.text_frame, "arch-quality team  |  2026-07", 13, color=GRAY, align=PP_ALIGN.CENTER)

# 2 - TOC
slide = new_slide("目录")
d = data["toc"]
for i, item in enumerate(d):
    y = Inches(1.3) + Inches(i * 0.9)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = NAVY; circ.line.fill.background()
    tf = circ.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = item["num"]; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Microsoft YaHei"
    tb = add_tb(slide, Inches(1.6), y + Inches(0.02), Inches(10), Inches(0.3))
    txt(tb.text_frame, item["title"], 18, True, NAVY)
    tb = add_tb(slide, Inches(1.6), y + Inches(0.32), Inches(10), Inches(0.25))
    txt(tb.text_frame, item["sub"], 12, color=GRAY)

# ===== Concept section (3-7) =====

# 3 - PDE solver intro
slide = new_slide("什么是 PDE 求解器？")
d = data["pde_solver_intro"]
tb = add_tb(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(3.5))
txt(tb.text_frame, d[0], 17, True, NAVY)
for i in range(1, len(d)):
    add_p(tb.text_frame, "\u2192  " + d[i], 16, color=DARK, sb=Pt(16))
add_box(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(1.2), fill=LIGHT_BLUE)
tb = add_tb(slide, Inches(1.0), Inches(5.0), Inches(11), Inches(0.8))
txt(tb.text_frame, "\u26a0\ufe0f  核心：控制离散误差在可接受范围", 20, True, NAVY)
add_p(tb.text_frame, "MMS 就是用来验证\u201c离散误差是否真的可控\u201d的系统方法", 14, color=DARK)

# 4 - Order comparison
slide = new_slide("收敛阶：理论阶 vs 观察阶")
oc = data["order_comparison"]
# Left card
add_box(slide, Inches(0.6), Inches(1.3), Inches(5.5), Inches(2.8), fill=LIGHT_GRAY, line=RGBColor(0xCC, 0xCC, 0xCC))
tb = add_tb(slide, Inches(0.8), Inches(1.4), Inches(5.1), Inches(0.4))
txt(tb.text_frame, oc["theory_title"], 18, True, NAVY)
tb = add_tb(slide, Inches(0.8), Inches(2.0), Inches(5.1), Inches(1.8))
txt(tb.text_frame, oc["theory_body"], 13, color=DARK, font="Consolas")
# Right card
add_box(slide, Inches(6.5), Inches(1.3), Inches(6.2), Inches(2.8), fill=LIGHT_GREEN, line=GREEN)
tb = add_tb(slide, Inches(6.7), Inches(1.4), Inches(5.8), Inches(0.4))
txt(tb.text_frame, oc["observed_title"], 18, True, GREEN)
tb = add_tb(slide, Inches(6.7), Inches(2.0), Inches(5.8), Inches(1.8))
txt(tb.text_frame, oc["observed_body"], 13, color=DARK, font="Consolas")
# Result cards
for i, (title, desc, accent, bg) in enumerate([
    ("\u2705  "+oc["result_pass"], "", GREEN, LIGHT_GREEN),
    ("\u274c  "+oc["result_fail"], "", RED, LIGHT_RED),
]):
    x = Inches(0.8) + Inches(i * 6.2)
    add_box(slide, x, Inches(4.5), Inches(5.8), Inches(1.0), fill=bg, line=accent)
    tb = add_tb(slide, x + Inches(0.3), Inches(4.6), Inches(5.2), Inches(0.6))
    txt(tb.text_frame, title, 17, True, accent)
add_box(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.8), fill=LIGHT_BLUE)
tb = add_tb(slide, Inches(1.0), Inches(6.0), Inches(11), Inches(0.5))
txt(tb.text_frame, "\U0001f4ca  " + oc["example"], 15, color=NAVY)

# 5 - Threshold basis
slide = new_slide("偏差阈值的工程依据")
d = data["threshold_basis"]
accents = [RED, ORANGE, GRAY]
bgs = [LIGHT_RED, LIGHT_ORANGE, LIGHT_GRAY]
for i, item in enumerate(d):
    x = Inches(0.6) + Inches(i * 4.2)
    ac = accents[i]; bg = bgs[i]
    add_box(slide, x, Inches(1.3), Inches(3.8), Inches(4.5), fill=WHITE, line=ac)
    # Title bar
    add_box(slide, x, Inches(1.3), Inches(3.8), Inches(0.6), fill=ac)
    tb = add_tb(slide, x + Inches(0.15), Inches(1.32), Inches(3.5), Inches(0.5))
    txt(tb.text_frame, item[0], 15, True, WHITE, align=PP_ALIGN.CENTER)
    # Reference
    tb = add_tb(slide, x + Inches(0.15), Inches(2.0), Inches(3.5), Inches(0.6))
    txt(tb.text_frame, item[1], 10, color=GRAY, align=PP_ALIGN.CENTER)
    # Range
    tb = add_tb(slide, x + Inches(0.15), Inches(2.7), Inches(3.5), Inches(0.4))
    txt(tb.text_frame, "偏差范围：" + item[2], 14, True, ac, align=PP_ALIGN.CENTER)
    # Formula / explanation
    tb = add_tb(slide, x + Inches(0.15), Inches(3.2), Inches(3.5), Inches(1.2))
    txt(tb.text_frame, item[3], 11, color=DARK, font="Consolas", align=PP_ALIGN.CENTER)
    # Bottom line
    tb = add_tb(slide, x + Inches(0.15), Inches(4.6), Inches(3.5), Inches(0.5))
    txt(tb.text_frame, item[4], 12, True, ac, align=PP_ALIGN.CENTER)
# Footer
add_box(slide, Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.7), fill=LIGHT_BLUE)
tb = add_tb(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5))
txt(tb.text_frame, "\U0001f4d6  " + data["threshold_basis"][0][1].split("\n")[0] + "  |  " + data["threshold_basis"][1][1].split("\n")[0] + "  |  " + data["threshold_basis"][2][1].split("\n")[0], 10, color=GRAY)

# 6 - Verify vs Validate
slide = new_slide("验证 vs 确认")
vv = data["verify_validate"]
# Verification card
add_box(slide, Inches(0.6), Inches(1.3), Inches(5.5), Inches(3.5), fill=LIGHT_GREEN, line=GREEN)
tb = add_tb(slide, Inches(0.8), Inches(1.5), Inches(5.1), Inches(0.5))
txt(tb.text_frame, "\u2705  " + vv["verification_title"], 22, True, GREEN)
tb = add_tb(slide, Inches(0.8), Inches(2.2), Inches(5.1), Inches(2.2))
txt(tb.text_frame, vv["verification_body"], 14, color=DARK)
# Validation card
add_box(slide, Inches(6.5), Inches(1.3), Inches(6.2), Inches(3.5), fill=LIGHT_BLUE, line=BLUE)
tb = add_tb(slide, Inches(6.7), Inches(1.5), Inches(5.8), Inches(0.5))
txt(tb.text_frame, "\U0001f4cb  " + vv["validation_title"], 22, True, BLUE)
tb = add_tb(slide, Inches(6.7), Inches(2.2), Inches(5.8), Inches(2.2))
txt(tb.text_frame, vv["validation_body"], 14, color=DARK)
# Divider arrow
add_box(slide, Inches(6.0), Inches(2.8), Inches(0.6), Pt(3), fill=RGBColor(0xBB, 0xBB, 0xBB))
# Footer
add_box(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.5), fill=LIGHT_ORANGE)
tb = add_tb(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(1.0))
txt(tb.text_frame, "\u26a0\ufe0f  " + vv["footer"], 14, True, ORANGE)
add_p(tb.text_frame, "本材料通篇讨论的是验证（Verification）。MMS 是验证的工具，不是确认的工具。", 13, color=DARK)

# 7 - Why systematic verification
slide = new_slide("为什么需要系统验证？")
d = data["why_system_verify"]
for i, item in enumerate(d):
    y = Inches(1.3) + Inches(i * 1.2)
    is_mms = (i == len(d) - 1)
    ac = GREEN if is_mms else GRAY
    bg = LIGHT_GREEN if is_mms else LIGHT_GRAY
    add_box(slide, Inches(0.8), y, Inches(11.5), Inches(0.95), fill=bg, line=ac)
    # Number
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.2), Inches(0.45), Inches(0.45))
    circ.fill.solid(); circ.fill.fore_color.rgb = GREEN if is_mms else BLUE; circ.line.fill.background()
    tf = circ.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i+1); r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Microsoft YaHei"
    # Title
    tb = add_tb(slide, Inches(1.7), y + Inches(0.1), Inches(3), Inches(0.35))
    txt(tb.text_frame, item[0], 16, True, ac)
    # Body
    tb = add_tb(slide, Inches(4.8), y + Inches(0.1), Inches(7), Inches(0.7))
    txt(tb.text_frame, item[1], 13, color=DARK)
    # Arrow to next
    if i < len(d) - 1:
        add_box(slide, Inches(6.0), y + Inches(0.95), Pt(2), Inches(0.3), fill=RGBColor(0xCC, 0xCC, 0xCC))
tb = add_tb(slide, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.4))
txt(tb.text_frame, "MMS 填补了传统方法无法系统验证 PDE 离散化实现正确性的空白。", 13, color=GRAY)

# 8 - Core problem
slide = new_slide("核心问题：代码的实现有没有 Bug？")
tb = add_tb(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.6))
txt(tb.text_frame, "写了一个 CFD/CAE 求解器 \u2014 怎么证明代码的实现没有 Bug？", 18, True, DARK)
tb = add_tb(slide, Inches(0.8), Inches(2.1), Inches(11.5), Inches(0.4))
txt(tb.text_frame, "工业仿真软件的 Bug 可能直接导致：", 14, color=GRAY)
d = data["disasters"]
for i, item in enumerate(d):
    x = Inches(0.8) + Inches(i * 4.1)
    c = {"red": RED, "orange": ORANGE}[item[2]]
    bg = {"red": LIGHT_RED, "orange": LIGHT_ORANGE}[item[2]]
    icons = ["\u2708\ufe0f", "\u2622\ufe0f", "\U0001f697"]
    add_box(slide, x, Inches(3.0), Inches(3.7), Inches(2.0), fill=bg, line=c)
    tb = add_tb(slide, x + Inches(0.2), Inches(3.2), Inches(3.3), Inches(0.5))
    txt(tb.text_frame, icons[i] + "  " + item[0], 17, True, c)
    tb = add_tb(slide, x + Inches(0.2), Inches(3.8), Inches(3.3), Inches(0.5))
    txt(tb.text_frame, item[1], 13, color=DARK)

# 9 - Traditional limits
slide = new_slide("传统验证方法的阿喀琉斯之踵")
d = [["方法", "局限"]] + data["traditional_limits"]
add_table(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(4.5), len(d), 2, d, [Inches(3), Inches(8.5)])
tb = add_tb(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5))
txt(tb.text_frame, "基准测试的假阴性：lid-driven cavity 通过 != 实际工程问题无 Bug。这个认知陷阱在工业界极为常见。", 12, color=GRAY)

# 10 - Birth of MMS
slide = new_slide("MMS 的诞生")
tb = add_tb(slide, Inches(0.8), Inches(1.1), Inches(11.5), Inches(0.5))
txt(tb.text_frame, 'Patrick Roache (1978): \u201c我们不需要精确解，我们可以自己制造一个。\u201d', 16, color=DARK)
tb.text_frame.paragraphs[0].font.italic = True
d = data["milestones"]
for i, item in enumerate(d):
    y = Inches(1.8) + Inches(i * 0.9)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.2), y + Inches(0.06), Inches(0.22), Inches(0.22))
    dot.fill.solid(); dot.fill.fore_color.rgb = GREEN if i == 4 else BLUE; dot.line.fill.background()
    if i < 4:
        add_box(slide, Inches(1.31), y + Inches(0.28), Pt(2), Inches(0.65), fill=RGBColor(0xCC, 0xCC, 0xCC))
    tb = add_tb(slide, Inches(1.7), y, Inches(1.2), Inches(0.35)); txt(tb.text_frame, item[0], 14, True, NAVY)
    tb = add_tb(slide, Inches(3.2), y, Inches(4), Inches(0.35)); txt(tb.text_frame, item[1], 14, True, DARK)
    tb = add_tb(slide, Inches(7.5), y, Inches(5), Inches(0.35)); txt(tb.text_frame, item[2], 12, color=GRAY)

# 11 - Reverse construction
slide = new_slide("核心思想：逆向构造")
add_box(slide, Inches(0.6), Inches(1.3), Inches(5.5), Inches(2.5), fill=LIGHT_GRAY)
tb = add_tb(slide, Inches(0.8), Inches(1.4), Inches(5), Inches(0.4))
txt(tb.text_frame, "X 传统方法（正向）", 16, True, RED)
tb = add_tb(slide, Inches(0.8), Inches(1.9), Inches(5), Inches(1.6))
txt(tb.text_frame, "已知精确解 + 简单 PDE\n         ->\n    求解器计算\n         ->\n    对比误差", 13, color=DARK)
add_p(tb.text_frame, "问题：精确解极其稀有，无法验证任意 PDE", 12, color=RED, sb=Pt(6))

add_box(slide, Inches(6.5), Inches(1.3), Inches(6.2), Inches(2.5), fill=LIGHT_GREEN)
tb = add_tb(slide, Inches(6.7), Inches(1.4), Inches(5.8), Inches(0.4))
txt(tb.text_frame, "O MMS 方法（逆向）", 16, True, GREEN)
tb = add_tb(slide, Inches(6.7), Inches(1.9), Inches(5.8), Inches(1.6))
txt(tb.text_frame, "选择 u* -> 推导源项 S\n         ->\n    求解带源项 PDE\n         ->\n    对比 u* 与 u_num", 13, color=DARK)
add_p(tb.text_frame, "关键：可以验证任意复杂的 PDE", 12, color=GREEN, sb=Pt(6))

add_box(slide, Inches(0.6), Inches(4.3), Inches(12.1), Inches(2.5), fill=LIGHT_BLUE)
tb = add_tb(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.0))
txt(tb.text_frame, "把验证问题从\u201c找精确解\u201d变成\u201c造精确解\u201d", 18, True, NAVY)
add_p(tb.text_frame, "后者的约束弱得多 \u2014 只要函数光滑且满足边界条件，任何函数都可作为制造解。", 13, color=DARK)
add_p(tb.text_frame, "这使得 MMS 能覆盖任意复杂的 PDE 系统：NS 方程、可压缩流、多相流、流固耦合...", 13, color=DARK)

# 12 - Five steps
slide = new_slide("MMS 五步法")
d = data["five_steps"]
for i, item in enumerate(d):
    y = Inches(1.3) + Inches(i * 1.1)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.08), Inches(0.42), Inches(0.42))
    circ.fill.solid(); circ.fill.fore_color.rgb = GREEN if i == 4 else BLUE; circ.line.fill.background()
    tf = circ.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = str(i + 1); r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Microsoft YaHei"
    if i < 4:
        add_box(slide, Inches(0.81), y + Inches(0.5), Pt(2), Inches(0.65), fill=RGBColor(0xCC, 0xCC, 0xCC))
    tb = add_tb(slide, Inches(1.3), y, Inches(2), Inches(0.35))
    txt(tb.text_frame, item[0] + ": " + item[1], 14, True, NAVY)
    tb = add_tb(slide, Inches(3.5), y, Inches(5), Inches(0.35))
    txt(tb.text_frame, item[2], 12, color=DARK, font="Consolas")
    tb = add_tb(slide, Inches(8.5), y, Inches(4), Inches(0.35))
    txt(tb.text_frame, item[3], 11, color=GRAY)

# 13 - Convergence order
slide = new_slide("MMS 的收敛阶判定")
tb = add_tb(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(0.8))
txt(tb.text_frame, "离散误差：||u_num - u_exact|| = C * h^p + 高阶小量", 17, True, NAVY)
add_p(tb.text_frame, "对数域：log E = log C + p * log h", 17, True, NAVY)

for i, (title, desc, accent, bg) in enumerate([
    ("O 验证通过", "p_obs ~ p_theory（如中心差分 -> 2）", GREEN, LIGHT_GREEN),
    ("X 代码存在 Bug", "p_obs 显著偏离理论阶", RED, LIGHT_RED),
]):
    x = Inches(0.8) + Inches(i * 6.2)
    add_box(slide, x, Inches(2.6), Inches(5.8), Inches(1.3), fill=bg, line=accent)
    tb = add_tb(slide, x + Inches(0.3), Inches(2.8), Inches(5.2), Inches(0.4))
    txt(tb.text_frame, title, 20, True, accent)
    tb = add_tb(slide, x + Inches(0.3), Inches(3.2), Inches(5.2), Inches(0.5))
    txt(tb.text_frame, desc, 13, color=DARK)

add_box(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.5), fill=LIGHT_ORANGE)
tb = add_tb(slide, Inches(1.0), Inches(4.5), Inches(11), Inches(2.0))
txt(tb.text_frame, "关于容差 +/- 0.1", 14, True, ORANGE)
add_p(tb.text_frame, "均匀网格 + 二阶格式：观察阶应极为接近 2.0（偏差 < 0.01）。若偏差 > 0.1 基本可判定代码有 Bug。", 12, color=DARK)
add_p(tb.text_frame, "建议容差：均匀网格 +/- 0.05 | 扭曲网格 +/- 0.15 | 非结构网格 +/- 0.2~0.3", 12, color=DARK)

# 14 - Bugs
slide = new_slide("MMS 能发现什么 Bug？")
d = data["bugs"]
add_table(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(4.5), len(d), 3, d, [Inches(4.5), Inches(3.5), Inches(3.5)])
tb = add_tb(slide, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5))
txt(tb.text_frame, "分区边界通量交换错误是并行 CFD 中最隐蔽的 Bug：串行跑对，并行跑错。MMS 的串并行对比是检测此类问题的金标准。", 12, color=GRAY)

# 15 - Standards
slide = new_slide("工业标准与合规")
d = data["standards"]
add_table(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(3.5), len(d), 3, d, [Inches(3.5), Inches(3.5), Inches(4.5)])
tb = add_tb(slide, Inches(0.8), Inches(5.3), Inches(11.5), Inches(0.6))
txt(tb.text_frame, "对于需要出口到美国市场的 CFD 软件，ASME V&V 20 合规是准入门槛。MMS 验证记录是审核的必备材料。", 13, True, NAVY)

# 16 - Application value
slide = new_slide("工业应用价值")
d = data["scenarios"]
icons = ["\U0001f52c", "\U0001f6a7", "\U0001f504", "\U0001f4ca", "\U0001f4cb"]
for i, item in enumerate(d):
    y = Inches(1.3) + Inches(i * 1.05)
    add_box(slide, Inches(0.8), y, Inches(11.5), Inches(0.85), fill=LIGHT_GRAY if i % 2 == 0 else WHITE)
    tb = add_tb(slide, Inches(1.2), y + Inches(0.05), Inches(10.8), Inches(0.3))
    txt(tb.text_frame, icons[i] + "  " + item[0], 16, True, NAVY)
    tb = add_tb(slide, Inches(1.2), y + Inches(0.4), Inches(10.8), Inches(0.35))
    txt(tb.text_frame, item[1], 13, color=DARK)

# 17 - Case setup
slide = new_slide("案例：稳态扩散 MMS")
add_box(slide, Inches(0.6), Inches(1.2), Inches(5.8), Inches(2.3), fill=LIGHT_GRAY)
tb = add_tb(slide, Inches(0.8), Inches(1.3), Inches(5.4), Inches(0.35))
txt(tb.text_frame, "问题设定", 14, True, NAVY)
tb = add_tb(slide, Inches(0.8), Inches(1.8), Inches(5.2), Inches(1.4))
txt(tb.text_frame, "方程：-div(DT*grad(T)) = S    DT=1", 12, color=DARK, font="Consolas")
add_p(tb.text_frame, "域：  [0,1] x [0,1]", 12, color=DARK, font="Consolas")
add_p(tb.text_frame, "BC：  T=0 on boundary", 12, color=DARK, font="Consolas")

add_box(slide, Inches(6.8), Inches(1.2), Inches(5.8), Inches(2.3), fill=LIGHT_GREEN)
tb = add_tb(slide, Inches(7.0), Inches(1.3), Inches(5.4), Inches(0.35))
txt(tb.text_frame, "制造解", 14, True, GREEN)
tb = add_tb(slide, Inches(7.0), Inches(1.8), Inches(5.4), Inches(0.7))
txt(tb.text_frame, "T*(x,y) = sin(pi*x) * sin(pi*y)", 15, True, DARK)
add_p(tb.text_frame, "光滑、边界为零（自动满足 Dirichlet BC）", 12, color=GRAY)

add_box(slide, Inches(0.6), Inches(4.0), Inches(12), Inches(1.3), fill=LIGHT_BLUE)
tb = add_tb(slide, Inches(0.8), Inches(4.1), Inches(11.5), Inches(0.35))
txt(tb.text_frame, "解析源项", 14, True, NAVY)
tb = add_tb(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.6))
txt(tb.text_frame, "S(x,y) = 2*pi^2 * sin(pi*x) * sin(pi*y)", 17, True, NAVY)

# 18 - Decision tree
slide = new_slide("求解器选型决策树")
tree_items = [
    ("需求：带源项的稳态扩散", True, None, DARK, LIGHT_GRAY2),
    ("候选：laplacianFoam", True, Inches(0.2), DARK, LIGHT_GRAY),
    ("  检查：是否支持 fvOptions？", False, None, DARK, None),
    ("X 不支持（被拒绝）", True, Inches(0.4), RED, LIGHT_RED),
    ("候选：scalarTransportFoam", True, Inches(0.2), DARK, LIGHT_GRAY),
    ("  策略：设 U=(0,0,0)，对流项归零", False, None, DARK, None),
    ("O 支持 fvOptions -> scalarCodedSource", True, Inches(0.4), GREEN, LIGHT_GREEN),
]
y = Inches(1.3)
for text, is_box, indent, c, bg in tree_items:
    if is_box:
        x = Inches(0.6) + (indent or Inches(0))
        w = Inches(11.5) - ((indent or Inches(0)).emu - Inches(0.6).emu if indent else 0)
        add_box(slide, x, y, w, Inches(0.42), fill=bg, line=RGBColor(0xCC, 0xCC, 0xCC))
        tb = add_tb(slide, x + Inches(0.1), y + Inches(0.04), w - Inches(0.2), Inches(0.34))
        txt(tb.text_frame, text, 13, True, c)
        y += Inches(0.5)
    else:
        tb = add_tb(slide, Inches(2.5), y, Inches(8), Inches(0.3))
        txt(tb.text_frame, text, 12, color=c)
        y += Inches(0.35)
tb = add_tb(slide, Inches(0.6), Inches(6.2), Inches(12), Inches(0.4))
txt(tb.text_frame, "核心原则：PDE 不变，求解器可换。laplacianFoam 不支持 fvOptions -> 换 scalarTransportFoam 即可。", 12, color=GRAY)

# 19 - Four pitfalls
slide = new_slide("四个坑")
d = data["pitfalls"]
accent = [RED, RED, ORANGE, ORANGE]
bg = [LIGHT_RED, LIGHT_RED, LIGHT_ORANGE, LIGHT_ORANGE]
for i, item in enumerate(d):
    x = Inches(0.6) + Inches(i % 2 * 6.2)
    y = Inches(1.3) + Inches(i // 2 * 2.8)
    add_box(slide, x, y, Inches(5.8), Inches(2.5), fill=bg[i], line=accent[i])
    tb = add_tb(slide, x + Inches(0.2), y + Inches(0.12), Inches(5.4), Inches(0.35))
    txt(tb.text_frame, item[0], 16, True, accent[i])
    tb = add_tb(slide, x + Inches(0.2), y + Inches(0.6), Inches(5.4), Inches(1.6))
    txt(tb.text_frame, item[1], 11, color=DARK, font="Consolas")

# 20 - Results
slide = new_slide("实测结果")
d = data["results"]
add_table(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.3), len(d), 6, d, [Inches(1.8), Inches(1.8), Inches(1.5), Inches(2.0), Inches(1.8), Inches(2.6)])
add_box(slide, Inches(4), Inches(4.0), Inches(5.3), Inches(1.5), fill=LIGHT_GREEN, line=GREEN)
tb = add_tb(slide, Inches(4.3), Inches(4.2), Inches(4.7), Inches(0.5))
txt(tb.text_frame, "O  p_obs = 2.001", 26, True, GREEN, align=PP_ALIGN.CENTER)
tb = add_tb(slide, Inches(4.3), Inches(4.8), Inches(4.7), Inches(0.4))
txt(tb.text_frame, "二阶收敛验证通过", 14, color=DARK, align=PP_ALIGN.CENTER)
tb = add_tb(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.5))
txt(tb.text_frame, "误差每加密一倍缩小约 4 倍 (4 = 2^2)，与二阶格式的理论特性完全吻合。实测偏差仅 0.001。", 12, color=GRAY)

# 21 - Extensions
slide = new_slide("扩展计划")
d = data["extensions"]
add_table(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.8), len(d), 6, d, [Inches(2.8), Inches(2.2), Inches(3.0), Inches(1.2), Inches(1.2), Inches(1.1)])
tb = add_tb(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(1.5))
txt(tb.text_frame, "扩展注意事项：", 14, True, NAVY)
add_p(tb.text_frame, "* simpleFoam：需设置 pRefCell，否则压力矩阵奇异", 12, color=DARK)
add_p(tb.text_frame, "* rhoCentralFoam：从亚声速光滑解开始，避免激波污染收敛阶", 12, color=DARK)
add_p(tb.text_frame, "* 多相流求解器：确保制造解满足 0 <= alpha <= 1", 12, color=DARK)

# 22 - One-click run
slide = new_slide("一键运行")
code_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.7),
         "$ cd docs/zh/mms\n$ python3 openfoam_laplacian_mms.py")
code_box(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(2.3),
         "[coarse] L2 error = 1.03e-03\n[medium] L2 error = 2.57e-04\n[fine]   L2 error = 6.43e-05\n\ncoarse -> medium: p_obs = 2.001  [PASS]\nmedium -> fine:   p_obs = 2.000  [PASS]\nOverall: [PASS]", bg=NAVY)
tb = add_tb(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.0))
txt(tb.text_frame, "所有配置自动生成，不需手动编辑 OpenFOAM case 文件。", 13, color=DARK)
add_p(tb.text_frame, "基于 OpenFOAM 2512 开发。v2306 以下版本将 scalarCodedSource 改为 codedSource。", 11, color=GRAY)

# 23 - Troubleshooting
slide = new_slide("排查失败五步法")
d = data["troubleshooting"]
for i, item in enumerate(d):
    y = Inches(1.3) + Inches(i * 1.1)
    accents = [RED, ORANGE, ORANGE, BLUE, BLUE]
    ac = accents[i]
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.05), Inches(0.45), Inches(0.45))
    circ.fill.solid(); circ.fill.fore_color.rgb = ac; circ.line.fill.background()
    tf = circ.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = item[0]; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Microsoft YaHei"
    if i < 4:
        add_box(slide, Inches(0.825), y + Inches(0.5), Pt(2), Inches(0.65), fill=RGBColor(0xDD, 0xDD, 0xDD))
    tb = add_tb(slide, Inches(1.4), y, Inches(3), Inches(0.35))
    txt(tb.text_frame, item[1], 16, True, ac)
    tb = add_tb(slide, Inches(4.5), y, Inches(8), Inches(0.7))
    txt(tb.text_frame, item[2], 12, color=DARK)

# 24 - CI/CD
slide = new_slide("CI/CD 集成策略")
d = data["ci_tiers"]
accent_map = {"green": GREEN, "orange": ORANGE, "red": RED}
icons_map = {"green": "\U0001f7e2", "orange": "\U0001f7e1", "red": "\U0001f534"}
for i, item in enumerate(d):
    x = Inches(0.6) + Inches(i * 4.2)
    ac = accent_map[item[2]]
    add_box(slide, x, Inches(1.3), Inches(3.8), Inches(3.8), fill=WHITE, line=ac)
    tb = add_tb(slide, x + Inches(0.3), Inches(1.5), Inches(3.2), Inches(0.5))
    txt(tb.text_frame, item[0].split("->")[0], 14, True, ac)
    tb = add_tb(slide, x + Inches(0.3), Inches(2.0), Inches(3.2), Inches(0.25))
    txt(tb.text_frame, "->", 12, color=GRAY, align=PP_ALIGN.CENTER)
    tb = add_tb(slide, x + Inches(0.3), Inches(2.2), Inches(3.2), Inches(0.3))
    parts = item[0].split("->")
    txt(tb.text_frame, parts[1].strip() if len(parts) > 1 else "", 14, True, ac)
    tb = add_tb(slide, x + Inches(0.3), Inches(2.7), Inches(3.2), Inches(2.0))
    txt(tb.text_frame, item[1], 11, color=DARK)
    if i < 2:
        tb = add_tb(slide, x + Inches(3.8), Inches(3.0), Inches(0.5), Inches(0.25))
        txt(tb.text_frame, "->", 14, color=GRAY)
tb = add_tb(slide, Inches(0.6), Inches(5.6), Inches(12), Inches(0.8))
txt(tb.text_frame, "关键原则：不让 MMS 成为开发的绊脚石。fast-check 秒级完成；全量回归放在 nightly/weekly。", 13, True, NAVY)

# 25 - Verification system
slide = new_slide("完整的验证体系")
d = data["verification_matrix"]
add_table(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(3.2), len(d), 3, d, [Inches(3), Inches(4), Inches(4.5)])
add_box(slide, Inches(4.5), Inches(5.3), Inches(4.3), Inches(0.8), fill=LIGHT_GREEN, line=GREEN)
tb = add_tb(slide, Inches(4.8), Inches(5.4), Inches(3.7), Inches(0.5))
txt(tb.text_frame, "MMS 验证通过是必要条件，不是充分条件", 13, True, GREEN, align=PP_ALIGN.CENTER)

# 26 - Roadmap
slide = new_slide("实施路线")
d = data["phases"]
for i, item in enumerate(d):
    y = Inches(1.1) + Inches(i * 1.0)
    colors = [RED, GREEN, GREEN, BLUE, BLUE, GRAY]
    c = colors[i]
    add_box(slide, Inches(0.6), y, Inches(1.4), Inches(0.38), fill=c)
    tb = add_tb(slide, Inches(0.7), y + Inches(0.02), Inches(1.2), Inches(0.33))
    txt(tb.text_frame, item[0], 10, True, WHITE, align=PP_ALIGN.CENTER)
    tb = add_tb(slide, Inches(2.2), y, Inches(2.5), Inches(0.38))
    txt(tb.text_frame, item[1], 14, True, DARK)
    tb = add_tb(slide, Inches(5.0), y, Inches(5.5), Inches(0.38))
    txt(tb.text_frame, item[2], 11, color=GRAY)
    # Status
    icons = ["\u2b50 P0"] + ["\u2705"] * 2 + ["\U0001f4c5"] * 3
    tb = add_tb(slide, Inches(11.0), y, Inches(1.5), Inches(0.38))
    txt(tb.text_frame, icons[i], 11, color=c, align=PP_ALIGN.RIGHT)

# 27 - Lifecycle
slide = new_slide("MMS 在软件生命周期中的定位")
items = ["需求规格", "架构设计", "代码实现", "单元测试", "集成测试", "交付"]
for i, item in enumerate(items):
    x = Inches(0.8) + Inches(i * 2.1)
    is_mms = (item == "单元测试")
    add_box(slide, x, Inches(1.5) + (Inches(0.4) if not is_mms else Inches(0)), Inches(1.8), Inches(0.8), fill=GREEN if is_mms else LIGHT_GRAY)
    tb = add_tb(slide, x + Inches(0.1), Inches(1.7) + (Inches(0.4) if not is_mms else Inches(0.1)), Inches(1.6), Inches(0.4))
    txt(tb.text_frame, item, 12, is_mms, DARK, align=PP_ALIGN.CENTER)
add_box(slide, Inches(0.8), Inches(2.8), Inches(12), Pt(3), fill=RGBColor(0xDD, 0xDD, 0xDD))
add_box(slide, Inches(3.8), Inches(3.1), Inches(4.5), Inches(1.8), fill=LIGHT_GREEN, line=GREEN)
tb = add_tb(slide, Inches(4.0), Inches(3.3), Inches(4.1), Inches(0.5))
txt(tb.text_frame, "\u25b2  MMS 验证", 20, True, GREEN, align=PP_ALIGN.CENTER)
tb = add_tb(slide, Inches(4.0), Inches(3.9), Inches(4.1), Inches(0.6))
txt(tb.text_frame, "PDE 求解器的\u201c单元测试\u201d\n验证（Verification）阶段", 13, color=DARK, align=PP_ALIGN.CENTER)
tb = add_tb(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.8))
txt(tb.text_frame, "MMS 属于验证（Verification），而非确认（Validation）。它位于单元测试和集成测试之间。", 12, color=GRAY)

# 28 - Summary
slide = new_slide("总结：三条核心信息")
d = data["summary_cards"]
accent = [GREEN, BLUE, ORANGE]
bg = [LIGHT_GREEN, LIGHT_BLUE, LIGHT_ORANGE]
for i, item in enumerate(d):
    y = Inches(1.3) + Inches(i * 1.8)
    ac = accent[i]
    add_box(slide, Inches(0.8), y, Inches(11.5), Inches(1.5), fill=bg[i], line=ac)
    tb = add_tb(slide, Inches(1.2), y + Inches(0.1), Inches(0.7), Inches(0.4))
    txt(tb.text_frame, item[0], 24, color=ac, align=PP_ALIGN.CENTER)
    tb = add_tb(slide, Inches(2.1), y + Inches(0.1), Inches(9.6), Inches(0.4))
    txt(tb.text_frame, item[1], 20, True, ac)
    tb = add_tb(slide, Inches(2.1), y + Inches(0.6), Inches(9.6), Inches(0.7))
    txt(tb.text_frame, item[2], 13, color=DARK)

# 29 - Q&A
slide = new_slide(dark=True)
tb = add_tb(slide, Inches(1), Inches(1.3), Inches(11), Inches(0.5))
txt(tb.text_frame, "推荐阅读", 20, True, WHITE, align=PP_ALIGN.CENTER)
d = data["references"]
for i, ref in enumerate(d):
    y = Inches(2.1) + Inches(i * 0.6)
    tb = add_tb(slide, Inches(1.5), y, Inches(10), Inches(0.45))
    txt(tb.text_frame, "\U0001f4d6  " + ref, 12, color=RGBColor(0xBB, 0xDD, 0xFF))
tb = add_tb(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5))
txt(tb.text_frame, "Thank You  |  Questions & Discussion", 22, True, WHITE, align=PP_ALIGN.CENTER)

prs.save(OUTPUT)
print("OK - PPT generated:", OUTPUT)
print("Total slides:", len(prs.slides))
