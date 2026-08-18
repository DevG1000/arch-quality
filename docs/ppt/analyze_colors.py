import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
import re

prs = Presentation(r'D:\opensource\arch-quality\docs\ppt\SKILL开发指南PPT.pptx')

def emu_to_inches(emu):
    if emu is None: return None
    return round(emu / 914400, 2)

def get_fill_color_from_xml(shape):
    """Extract fill color from shape XML"""
    spPr = shape._element.findall('.//' + qn('a:spPr'))
    results = {}
    for pr in spPr:
        pr_xml = pr.xml
        # Check solid fill
        m = re.search(r'<a:solidFill>.*?<a:srgbClr val="([0-9A-Fa-f]+)"/>', pr_xml)
        if m:
            results['fill'] = m.group(1).lower()
        # Check line fill
        m2 = re.search(r'<a:ln.*?<a:solidFill>.*?<a:srgbClr val="([0-9A-Fa-f]+)"/>', pr_xml, re.DOTALL)
        if m2:
            results['line'] = m2.group(1).lower()
        # Check if rounded
        if 'round' in pr_xml:
            results['rounded'] = True
    return results

print("=== COLOR & FONT DETAILED ANALYSIS ===")
print()

# Phase 1: Check title bar fills (Rounded Rectangle 1 - the top bar)
print("--- TITLE BAR (Shape 1) Colors ---")
for idx, slide in enumerate(prs.slides):
    n = idx + 1
    for shape in slide.shapes:
        if shape.name == 'Rounded Rectangle 1':
            fill_info = get_fill_color_from_xml(shape)
            print(f"  Slide {n}: {fill_info}")

print()
print("--- ICON BACKGROUND (Rounded Rectangle 2) Colors ---")
for idx, slide in enumerate(prs.slides):
    n = idx + 1
    for shape in slide.shapes:
        if shape.name == 'Rounded Rectangle 2':
            fill_info = get_fill_color_from_xml(shape)
            print(f"  Slide {n}: {fill_info}")

print()
print("--- CARD FILLS (Rounded Rectangle 7, 9) Colors ---")
for idx, slide in enumerate(prs.slides):
    n = idx + 1
    for shape in slide.shapes:
        if shape.name in ('Rounded Rectangle 7', 'Rounded Rectangle 9'):
            fill_info = get_fill_color_from_xml(shape)
            print(f"  Slide {n} {shape.name}: {fill_info}")

print()
print("--- ALL FILL COLORS SUMMARY ---")
color_summary = {}
for idx, slide in enumerate(prs.slides):
    n = idx + 1
    for shape in slide.shapes:
        fi = get_fill_color_from_xml(shape)
        fill = fi.get('fill')
        if fill:
            key = f"#{fill}"
            if key not in color_summary:
                color_summary[key] = []
            color_summary[key].append((n, shape.name))

for color, locations in sorted(color_summary.items()):
    print(f"  {color}: {locations}")

print()
print("--- FONT SIZES & COLORS (all text) ---")
font_stats = {}
for idx, slide in enumerate(prs.slides):
    n = idx + 1
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if not txt:
                    continue
                for run in para.runs:
                    sz = run.font.size
                    bold = run.font.bold
                    try:
                        color = str(run.font.color.rgb).lower()
                    except:
                        color = 'none'
                    font_name = run.font.name
                    
                    if font_name:
                        if font_name not in font_stats:
                            font_stats[font_name] = {'count': 0, 'slides': set(), 'sizes': set(), 'colors': set()}
                        font_stats[font_name]['count'] += 1
                        font_stats[font_name]['slides'].add(n)
                        if sz:
                            font_stats[font_name]['sizes'].add(round(sz/12700, 1))
                        font_stats[font_name]['colors'].add(color)

for font, stats in sorted(font_stats.items()):
    sz_str = ', '.join(str(s) for s in sorted(stats['sizes'])) if stats['sizes'] else 'inherited'
    col_str = ', '.join(sorted(stats['colors']))
    print(f"  Font '{font}': used {stats['count']}x on slides {sorted(stats['slides'])}")
    print(f"    Sizes: {sz_str} pt")
    print(f"    Colors: {col_str}")

print()
print("--- SLIDE 1 & 22 SHAPE DETAILS ---")
for target in [1, 22]:
    slide = prs.slides[target - 1]
    print(f"  Slide {target}:")
    for shape in slide.shapes:
        print(f"    {shape.name}: left={emu_to_inches(shape.left)} top={emu_to_inches(shape.top)} w={emu_to_inches(shape.width)} h={emu_to_inches(shape.height)}")
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    align = para.alignment
                    for run in para.runs:
                        sz = run.font.size
                        sz_pt = round(sz/12700, 1) if sz else 'inherited'
                        print(f"      '{txt[:60]}' align={align} font={run.font.name} size={sz_pt}pt bold={run.font.bold}")
                        break
