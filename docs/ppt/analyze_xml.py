import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

prs = Presentation(r'D:\opensource\arch-quality\docs\ppt\SKILL开发指南PPT.pptx')

# Look at shape XML for slide 2, shape 1 (title bar) to understand color structure
slide = prs.slides[1]  # Slide 2
for shape in slide.shapes:
    if shape.name == 'Rounded Rectangle 1':
        xml_str = etree.tostring(shape._element, pretty_print=True).decode('utf-8')
        print("=== Rounded Rectangle 1 (Slide 2) XML ===")
        print(xml_str[:2000])
        break

print()

# Also check shape 2 (icon background)
for shape in slide.shapes:
    if shape.name == 'Rounded Rectangle 2':
        xml_str = etree.tostring(shape._element, pretty_print=True).decode('utf-8')
        print("=== Rounded Rectangle 2 (Slide 2) XML ===")
        print(xml_str[:2000])
        break

print()

# Check shape 7 (card)
for shape in slide.shapes:
    if shape.name == 'Rounded Rectangle 7':
        xml_str = etree.tostring(shape._element, pretty_print=True).decode('utf-8')
        print("=== Rounded Rectangle 7 (Slide 2) XML ===")
        print(xml_str[:2000])
        break

print()

# Check slide 1 - the main rectangle (shape 4)
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.name in ('Rounded Rectangle 4',):
        xml_str = etree.tostring(shape._element, pretty_print=True).decode('utf-8')
        print(f"=== {shape.name} (Slide 1) XML ===")
        print(xml_str[:2000])
        break
