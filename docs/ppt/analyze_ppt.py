from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

prs = Presentation(r'D:\opensource\arch-quality\docs\ppt\SKILL开发指南PPT.pptx')

slide_count = len(prs.slides)
print(f"Slide count: {slide_count}")
print(f"Slide width: {prs.slide_width}  height: {prs.slide_height}")
print()

def emu_to_inches(emu):
    if emu is None:
        return None
    return round(emu / 914400, 2)

# 1. Collect all slide data
slides_data = []
for idx, slide in enumerate(prs.slides):
    n = idx + 1
    info = {"num": n, "shapes": [], "bg_color": None}
    bg = slide.background
    fill = bg.fill
    if fill.type is not None:
        try:
            if fill.type == 1:
                c = fill.fore_color.rgb
                info["bg_color"] = str(c)
        except:
            pass
    for shape in slide.shapes:
        s = {
            "name": shape.name,
            "type": str(shape.shape_type),
            "left": shape.left,
            "top": shape.top,
            "w": shape.width,
            "h": shape.height,
            "texts": []
        }
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                txt = p.text.strip()
                if txt:
                    for run in p.runs:
                        sz = run.font.size
                        bold = run.font.bold
                        color = None
                        try:
                            color = str(run.font.color.rgb)
                        except:
                            pass
                        font_name = run.font.name
                        s["texts"].append({
                            "text": txt[:80],
                            "size": sz,
                            "bold": bold,
                            "color": color,
                            "font": font_name,
                            "align": p.alignment
                        })
                        break
        info["shapes"].append(s)
    slides_data.append(info)

# ============ CHECK 1: Title bar icons ============
print("=" * 70)
print("CHECK 1: TITLE BAR ICONS")
print("=" * 70)
# Icons are typically small squares (left < ~1 inch, width < ~1 inch) at left side of slide
# We look for shapes containing single Unicode-like characters or very small shapes

icon_candidates = []
missing_icon_slides = []
icon_details = []

for sd in slides_data:
    n = sd["num"]
    found_icon = False
    for sh in sd["shapes"]:
        left_in = emu_to_inches(sh["left"])
        top_in = emu_to_inches(sh["top"])
        w_in = emu_to_inches(sh["w"])
        h_in = emu_to_inches(sh["h"])
        # Look for small shapes on left side
        if left_in is not None and left_in < 2.0 and w_in is not None and w_in < 1.5 and h_in is not None and h_in < 1.0:
            for t in sh["texts"]:
                txt = t["text"]
                if len(txt) <= 4 and not txt.startswith("Slide"):
                    found_icon = True
                    icon_details.append((n, sh["name"], txt, left_in, top_in, w_in, h_in))
                    break
    if found_icon:
        print(f"  Slide {n}: HAS ICON")
    else:
        missing_icon_slides.append(n)
        print(f"  Slide {n}: MISSING ICON")

print()
print(f"Slides WITH icons: {[sd['num'] for sd in slides_data if sd['num'] not in missing_icon_slides]}")
print(f"Slides MISSING icons: {missing_icon_slides}")

# Check overlap: icon vs title text
print()
print("--- Icon positioning details ---")
for n, name, txt, left, top, w, h in icon_details:
    safe_txt = txt.encode('ascii', 'replace').decode('ascii')
    print(f"  Slide {n}: icon='{safe_txt}' left={left}in top={top}in size={w}x{h}in")

# Check if any icon overlaps with title (title is usually near top-left too)
print()
overlap_issues = []
for sd in slides_data:
    n = sd["num"]
    icons_f = [(s["left"], s["w"]) for s in sd["shapes"] if s["name"].startswith("Icon") or s["name"].startswith("icon")]
    # Better: find shapes that look like icons and shapes that look like titles
    icon_rects = []
    title_rects = []
    for sh in sd["shapes"]:
        left_in = emu_to_inches(sh["left"])
        top_in = emu_to_inches(sh["top"])
        w_in = emu_to_inches(sh["w"])
        h_in = emu_to_inches(sh["h"])
        for t in sh["texts"]:
            txt = t["text"]
            # Identify icon: short single string, small width, left side
            if left_in is not None and left_in < 2.0 and w_in is not None and w_in < 1.2 and len(txt) <= 4:
                icon_rects.append((left_in, top_in, left_in + w_in, top_in + h_in, txt))
            # Identify title: first or large bold text at top
            if t.get("bold") and t.get("size") and t["size"] >= Pt(18).emu:
                title_rects.append((left_in, top_in, left_in + w_in, top_in + h_in, txt))
    if icon_rects and title_rects:
        for ir in icon_rects:
            for tr in title_rects:
                # Check horizontal overlap: icons are left of title, so we check if icon right > title left
                if ir[0] is not None and tr[0] is not None and ir[2] is not None and tr[0] is not None:
                    if ir[2] > tr[0] - 0.1:  # within 0.1 inch of each other
                        overlap_issues.append((n, ir[4], tr[4]))
                        sic = ir[4].encode('ascii', 'replace').decode('ascii')
                        stc = tr[4].encode('ascii', 'replace').decode('ascii')
                        print(f"  Slide {n}: Icon '{sic}' may overlap with title '{stc}' (icon_right={ir[2]}in, title_left={tr[0]}in)")

if not overlap_issues:
    print("  No icon-title overlap detected.")
else:
    for n, ic, ti in overlap_issues:
        sic = ic.encode('ascii', 'replace').decode('ascii')
        sti = ti.encode('ascii', 'replace').decode('ascii')
        print(f"  Slide {n}: icon '{sic}' overlaps/near title '{sti}'")

# ============ CHECK 2: Visual consistency ============
print()
print("=" * 70)
print("CHECK 2: VISUAL CONSISTENCY")
print("=" * 70)

# 2a: Dark background
print()
print("--- Dark background (#1a1a2e) ---")
non_dark_slides = []
for sd in slides_data:
    n = sd["num"]
    c = sd["bg_color"]
    if c and c.upper() == "1A1A2E":
        pass  # ok
    else:
        non_dark_slides.append((n, c))
        print(f"  Slide {n}: bg_color={c} (expected 1a1a2e)")

if not non_dark_slides:
    print("  PASS: All slides have #1a1a2e background (or background is inherited)")
else:
    print(f"  ISSUE: {len(non_dark_slides)} slides with different background")

# 2b: Card elements (dark rounded rectangles)
print()
print("--- Card elements (dark rounded rectangles) ---")
card_slides = []
no_card_slides = []
for sd in slides_data:
    n = sd["num"]
    has_card = False
    for sh in sd["shapes"]:
        # Cards are typically large rectangles with dark fill
        w_in = emu_to_inches(sh["w"])
        h_in = emu_to_inches(sh["h"])
        if w_in is not None and h_in is not None and w_in > 2.0 and h_in > 1.0 and "Rect" in sh["type"]:
            # Check if it has a dark fill
            has_card = True
            break
    if has_card:
        card_slides.append(n)
    else:
        no_card_slides.append(n)

print(f"  Slides WITH card shapes: {card_slides}")
print(f"  Slides WITHOUT card shapes: {no_card_slides}")

# 2c: Color scheme
print()
print("--- Color scheme (ACCENT #00d4ff, GREEN #4ecb71, RED #ff6b6b) ---")
color_usage = {"00d4ff": 0, "4ecb71": 0, "ff6b6b": 0, "other": {}}
for sd in slides_data:
    for sh in sd["shapes"]:
        for t in sh["texts"]:
            c = t["color"]
            if c:
                cu = c.upper()
                if cu == "00D4FF":
                    color_usage["00d4ff"] += 1
                elif cu == "4ECB71":
                    color_usage["4ecb71"] += 1
                elif cu == "FF6B6B":
                    color_usage["ff6b6b"] += 1
                else:
                    if cu not in color_usage["other"]:
                        color_usage["other"][cu] = []
                    color_usage["other"][cu].append(sd["num"])

print(f"  ACCENT #00d4ff used: {color_usage['00d4ff']} times")
print(f"  GREEN #4ecb71 used: {color_usage['4ecb71']} times")
print(f"  RED #ff6b6b used: {color_usage['ff6b6b']} times")
other_colors = {k: len(v) for k, v in color_usage["other"].items()}
print(f"  Other colors found: {other_colors}")

# ============ CHECK 3: Spacing & Layout ============
print()
print("=" * 70)
print("CHECK 3: SPACING & LAYOUT")
print("=" * 70)

print()
print("--- Content overflow check ---")
slide_w_in = emu_to_inches(prs.slide_width)
slide_h_in = emu_to_inches(prs.slide_height)
print(f"  Slide dimensions: {slide_w_in} x {slide_h_in} inches")

overflow_slides = []
for sd in slides_data:
    n = sd["num"]
    for sh in sd["shapes"]:
        left_in = emu_to_inches(sh["left"])
        top_in = emu_to_inches(sh["top"])
        right_in = emu_to_inches(sh["left"] + sh["w"])
        bottom_in = emu_to_inches(sh["top"] + sh["h"])
        if left_in is not None and left_in < -0.1:
            overflow_slides.append((n, f"left overflow: {left_in}"))
        elif right_in is not None and right_in > slide_w_in + 0.1:
            overflow_slides.append((n, f"right overflow: right={right_in} vs slide={slide_w_in}"))
        if top_in is not None and top_in < -0.1:
            overflow_slides.append((n, f"top overflow: {top_in}"))
        elif bottom_in is not None and bottom_in > slide_h_in + 0.1:
            overflow_slides.append((n, f"bottom overflow: bottom={bottom_in} vs slide={slide_h_in}"))

if overflow_slides:
    print(f"  OVERFLOW ISSUES found:")
    for n, msg in overflow_slides:
        print(f"    Slide {n}: {msg}")
else:
    print("  No overflow issues detected.")

print()
print("--- Left/right balance check ---")
# Check if text/content is concentrated on one side
balance_issues = []
for sd in slides_data:
    n = sd["num"]
    left_items = 0
    right_items = 0
    for sh in sd["shapes"]:
        left_in = emu_to_inches(sh["left"])
        w_in = emu_to_inches(sh["w"])
        if left_in is not None and w_in is not None:
            center = left_in + w_in / 2
            mid = slide_w_in / 2
            if center < mid - 0.5:
                left_items += 1
            elif center > mid + 0.5:
                right_items += 1
    if left_items > 0 and right_items == 0:
        balance_issues.append((n, f"all content on left side ({left_items} items)"))
    elif right_items > 0 and left_items == 0:
        balance_issues.append((n, f"all content on right side ({right_items} items)"))

if balance_issues:
    print(f"  UNEVEN BALANCE on slides:")
    for n, msg in balance_issues:
        print(f"    Slide {n}: {msg}")
else:
    print("  No extreme left/right imbalance detected.")

# ============ CHECK 4: Professionalism ============
print()
print("=" * 70)
print("CHECK 4: PROFESSIONALISM")
print("=" * 70)

print()
print("--- Font consistency ---")
# We look at all text runs and check font names
font_usage = {}
for sd in slides_data:
    for sh in sd["shapes"]:
        for t in sh["texts"]:
            fn = t["font"]
            if fn:
                if fn not in font_usage:
                    font_usage[fn] = set()
                font_usage[fn].add(sd["num"])

for fn, slides_list in sorted(font_usage.items()):
    print(f"  Font '{fn}' used on slides: {sorted(slides_list)}")

# Check for YaHei and FangSong
print()
yahei_slides = font_usage.get("Microsoft YaHei", set())
fangsong_slides = font_usage.get("FangSong", set())
print(f"  Microsoft YaHei used on: {sorted(yahei_slides)}")
print(f"  FangSong used on: {sorted(fangsong_slides)}")

# 4b: Visual hierarchy
print()
print("--- Visual hierarchy ---")
hierarchy_notes = []
for sd in slides_data:
    n = sd["num"]
    sizes_seen = set()
    for sh in sd["shapes"]:
        for t in sh["texts"]:
            if t["size"]:
                sz_pt = round(t["size"] / 12700, 1)
                sizes_seen.add(sz_pt)
    if len(sizes_seen) > 0:
        sorted_sizes = sorted(sizes_seen)
        hierarchy_notes.append((n, sorted_sizes))
        if len(sorted_sizes) >= 2 and sorted_sizes[-1] >= sorted_sizes[0] * 1.5:
            pass  # good hierarchy
        else:
            print(f"  Slide {n}: limited size variation: {sorted_sizes}")

print(f"  Hierarchy analysis done for all slides.")

# 4c: Slide 1 and 22 (title/end) layout
print()
print("--- Slide 1 and {last} (title/end) layout ---")
last_n = slide_count
for target in [1, last_n]:
    sd = slides_data[target - 1]
    print(f"  Slide {target}:")
    for sh in sd["shapes"]:
        left_in = emu_to_inches(sh["left"])
        top_in = emu_to_inches(sh["top"])
        w_in = emu_to_inches(sh["w"])
        h_in = emu_to_inches(sh["h"])
        # For title slides, text should be roughly centered
        for t in sh["texts"]:
            txt = t["text"][:60]
            align = t["align"]
            sz = t["size"]
            if sz:
                sz_pt = round(sz / 12700, 1)
            else:
                sz_pt = None
            safe_txt2 = txt.encode('ascii', 'replace').decode('ascii')
            print(f"    '{safe_txt2}' | pos=({left_in},{top_in}) size=({w_in}x{h_in}) | font={t['font']} {sz_pt}pt bold={t['bold']} align={align}")

print()
print("=" * 70)
print("FINAL REPORT")
print("=" * 70)
