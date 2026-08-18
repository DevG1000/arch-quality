import sys
import io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
import re

prs = Presentation(r'D:\opensource\arch-quality\docs\ppt\SKILL开发指南PPT.pptx')

def emu_to_inches(emu):
    if emu is None: return None
    return round(emu / 914400, 2)

print(f"Slide count: {len(prs.slides)}")
print(f"Slide size: {emu_to_inches(prs.slide_width)}\" x {emu_to_inches(prs.slide_height)}\"")
print()

# Deep dive into shapes per slide
for idx, slide in enumerate(prs.slides):
    n = idx + 1
    print(f"=== Slide {n} ===")
    
    # Check background XML for dark color
    bg_xml = slide.background._element.xml[:500]
    if '1a1a2e' in bg_xml.lower():
        print(f"  Background: HAS #1a1a2e (in XML)")
    else:
        print(f"  Background: checking XML...")
        print(f"    {bg_xml}")
    
    for shape in slide.shapes:
        left_in = emu_to_inches(shape.left)
        top_in = emu_to_inches(shape.top)
        w_in = emu_to_inches(shape.width)
        h_in = emu_to_inches(shape.height)
        
        # Get shape fill from XML
        spPr = shape._element.findall('.//' + qn('a:spPr'))
        fill_info = ""
        for pr in spPr:
            pr_xml = pr.xml
            if '1a1a2e' in pr_xml.lower():
                fill_info = "fill=#1a1a2e"
            elif '00d4ff' in pr_xml.lower() or '4ecb71' in pr_xml.lower() or 'ff6b6b' in pr_xml.lower():
                colors = re.findall(r'(00d4ff|4ecb71|ff6b6b)', pr_xml.lower())
                fill_info = f"fill=#{','.join(colors)}"
            elif 'srgb' in pr_xml.lower():
                m = re.search(r'srgbClr val="([0-9A-Fa-f]+)"', pr_xml)
                if m:
                    fill_info = f"fill=#{m.group(1)}"
        
        # Check for rounded rectangle
        is_rounded = False
        for pr in spPr:
            if 'round' in pr.xml.lower():
                is_rounded = True
        
        print(f"  Shape: {shape.name} type={shape.shape_type} pos=({left_in},{top_in}) size=({w_in}x{h_in}) {fill_info}{' [ROUNDED]' if is_rounded else ''}")
        
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                txt = para.text.strip()
                if txt:
                    align = para.alignment
                    for run in para.runs:
                        sz = run.font.size
                        bold = run.font.bold
                        font_name = run.font.name
                        try:
                            color = str(run.font.color.rgb)
                        except:
                            color = None
                        print(f"    Text: \"{txt[:80]}\" font={font_name} size={sz} bold={bold} color={color} align={align}")
                        break
    print()
