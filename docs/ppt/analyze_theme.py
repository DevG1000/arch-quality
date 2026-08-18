import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from lxml import etree

prs = Presentation(r'D:\opensource\arch-quality\docs\ppt\SKILL开发指南PPT.pptx')

# Check theme for fonts - iterate over package parts
import zipfile
import io as python_io
with zipfile.ZipFile(r'D:\opensource\arch-quality\docs\ppt\SKILL开发指南PPT.pptx', 'r') as z:
    for name in z.namelist():
        if 'theme' in name.lower():
            print(f"Part: {name}")
            with z.open(name) as f:
                xml_str = f.read().decode('utf-8')
            import re
            major = re.search(r'<a:majorFont>(.*?)</a:majorFont>', xml_str, re.DOTALL)
            minor = re.search(r'<a:minorFont>(.*?)</a:minorFont>', xml_str, re.DOTALL)
            if major:
                ea = re.search(r'ea typeface="([^"]+)"', major.group(1))
                lat = re.search(r'<a:latin typeface="([^"]+)"/>', major.group(1))
                print(f"  Major (headings): latin={lat.group(1) if lat else 'N/A'}, ea={ea.group(1) if ea else 'N/A'}")
            if minor:
                ea = re.search(r'ea typeface="([^"]+)"', minor.group(1))
                lat = re.search(r'<a:latin typeface="([^"]+)"/>', minor.group(1))
                print(f"  Minor (body): latin={lat.group(1) if lat else 'N/A'}, ea={ea.group(1) if ea else 'N/A'}")
            
            # Extract color scheme
            if 'clrScheme' in xml_str:
                colors = re.findall(r'<a:([a-z]+)\s+.*?val="([^"]+)"', xml_str)
                for name2, val in colors[:25]:
                    print(f"  Color: {name2} = {val}")
            break

# Now let's check the slide layouts for font info
print()
print("--- Slide Layouts ---")
for layout in prs.slide_layouts:
    print(f"Layout: {layout.name}")
    for ph in layout.placeholders:
        print(f"  Placeholder idx={ph.placeholder_format.idx} type={ph.placeholder_format.type} name={ph.name}")

# Check slide 1 text runs more carefully
print()
print("=== Slide 1: text runs detail ===")
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            txt = para.text.strip()
            if txt:
                for run in para.runs:
                    rPr = run._r.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                    font_xml = ''
                    for rp in rPr:
                        font_xml = etree.tostring(rp, pretty_print=True).decode('utf-8')
                    print(f"  Text: '{txt[:60]}'")
                    if font_xml:
                        print(f"    rPr: {font_xml[:300]}")
                    else:
                        print(f"    rPr: (none -> theme inherited)")

print()
print("=== Slide 22: text runs detail ===")
slide22 = prs.slides[21]
for shape in slide22.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            txt = para.text.strip()
            if txt:
                for run in para.runs:
                    rPr = run._r.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                    font_xml = ''
                    for rp in rPr:
                        font_xml = etree.tostring(rp, pretty_print=True).decode('utf-8')
                    print(f"  Text: '{txt[:60]}'")
                    if font_xml:
                        print(f"    rPr: {font_xml[:300]}")
                    else:
                        print(f"    rPr: (none -> theme inherited)")

# Check content slides for font specification
print()
print("=== Slide 2: text runs detail ===")
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            txt = para.text.strip()
            if txt:
                for run in para.runs:
                    rPr = run._r.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                    for rp in rPr:
                        sz = rp.get('sz')
                        b = rp.get('b')
                        lang = rp.get('lang')
                        latin = rp.find('{http://schemas.openxmlformats.org/drawingml/2006/main}latin')
                        ea = rp.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
                        solidFill = rp.find('{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
                        color = None
                        if solidFill is not None:
                            sc = solidFill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                            if sc is not None:
                                color = sc.get('val')
                        print(f"    '{txt[:40]}' sz={sz} b={b} lang={lang} latin={(latin.get('typeface') if latin is not None else None)} ea={(ea.get('typeface') if ea is not None else None)} color={color}")
                    else:
                        print(f"    '{txt[:40]}' -> NO rPr (inherited)")
