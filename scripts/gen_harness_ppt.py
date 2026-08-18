"""生成 Harness 工程科普 PPT。

参照 docs/ppt/SKILL开发指南PPT 的视觉风格（Blank 版式 + TEXT_BOX 自绘），
16:9 尺寸。16 页，概念:实操 = 6:4。
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---- 配色（对齐现有 PPT 风格）----
NAVY = RGBColor(0x1F, 0x3A, 0x5F)       # 深蓝标题
DARK = RGBColor(0x33, 0x33, 0x33)       # 深灰正文
GRAY = RGBColor(0x66, 0x66, 0x66)       # 浅灰副题
ACCENT = RGBColor(0xC0, 0x50, 0x4D)     # 强调红
LIGHT = RGBColor(0xF2, 0xF2, 0xF2)      # 浅灰背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)      # 通过绿

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUT_PATH = os.path.join(os.path.dirname(__file__), "..", "docs", "ppt",
                        "Harness工程科普.pptx")


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color=DARK,
                align=PP_ALIGN.LEFT, mono=False, line_spacing=1.0):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        if mono:
            run.font.name = "Consolas"
    return tb


def add_title(slide, chapter, title, subtitle=None):
    """章节页头：左上角章节标签 + 大标题 + 副题"""
    add_textbox(slide, 0.6, 0.4, 12, 0.4, chapter, size=13, bold=True, color=ACCENT)
    add_textbox(slide, 0.6, 0.8, 12, 0.9, title, size=32, bold=True, color=NAVY)
    if subtitle:
        add_textbox(slide, 0.6, 1.65, 12, 0.5, subtitle, size=15, color=GRAY)


def add_footer(slide, page_no):
    add_textbox(slide, 12.3, 7.05, 0.8, 0.3, str(page_no), size=11, color=GRAY,
                align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, data, header=True):
    """简单表格：data = [[...]]，第一行为表头"""
    rows, cols = len(data), len(data[0])
    from pptx.util import Inches as _In
    shape = slide.shapes.add_table(rows, cols, _In(x), _In(y), _In(w),
                                   _In(0.25 * rows)).table
    for ci in range(cols):
        shape.columns[ci].width = _In(w / cols)
    for ri in range(rows):
        for ci in range(cols):
            cell = shape.cell(ri, ci)
            cell.text = str(data[ri][ci])
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = DARK
                    if header and ri == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
            if header and ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
    return shape


def add_code(slide, x, y, w, h, text, size=12):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.name = "Consolas"
        run.font.color.rgb = NAVY


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]  # Blank

    # ============ 页1 封面 ============
    s = prs.slides.add_slide(blank)
    add_textbox(s, 1.5, 2.0, 10, 1.2, "Harness 工程科普", size=48, bold=True,
                color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s, 1.5, 3.2, 10, 0.7, "模型之外，可靠之道", size=24,
                color=GRAY, align=PP_ALIGN.CENTER)
    add_textbox(s, 1.5, 4.2, 10, 0.5,
                "Agent = Model + [上下文+工具+约束+验证+纠正]", size=18,
                color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, 1.5, 5.2, 10, 0.4,
                "基于 arch-quality 项目实操 · 2026-08", size=14,
                color=GRAY, align=PP_ALIGN.CENTER)
    add_footer(s, 1)

    # ============ 页2 目录 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "CONTENTS", "目录", "Why → What → How → Practice")
    items = [
        ("01", "为什么需要 Harness", "裸 LLM 的脆弱点与产品鸿沟"),
        ("02", "Harness 是什么", "核心公式与五要素保障层"),
        ("03", "怎么搭 Harness", "基于 opencode 的三层验证架构"),
        ("04", "效果与反思", "验证结果、踩过的坑、未来边界"),
    ]
    y = 2.4
    for num, t, sub in items:
        add_textbox(s, 1.0, y, 1.2, 0.6, num, size=26, bold=True, color=ACCENT)
        add_textbox(s, 2.2, y, 9, 0.5, t, size=20, bold=True, color=NAVY)
        add_textbox(s, 2.2, y + 0.5, 9, 0.4, sub, size=13, color=GRAY)
        y += 1.1
    add_footer(s, 2)

    # ============ 页3 裸 LLM 的脆弱点 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "01 为什么需要 Harness", "裸 LLM 的三个脆弱点",
              "ReAct 循环有效，但不够可靠")
    add_textbox(s, 1.0, 2.3, 11, 0.5, "1. 幻觉（Hallucination）", size=20, bold=True, color=NAVY)
    add_textbox(s, 1.0, 2.8, 11, 0.5, "编造不存在的工具、参数，甚至编造执行结果", size=15, color=DARK)
    add_textbox(s, 1.0, 3.5, 11, 0.5, "2. 选错工具（Wrong Tool）", size=20, bold=True, color=NAVY)
    add_textbox(s, 1.0, 4.0, 11, 0.5, "该调评估命令时去读源码，该读文件时去跑命令", size=15, color=DARK)
    add_textbox(s, 1.0, 4.7, 11, 0.5, "3. 无法自愈（No Recovery）", size=20, bold=True, color=NAVY)
    add_textbox(s, 1.0, 5.2, 11, 0.5, "遇到错误反复重试同一步，陷入死循环", size=15, color=DARK)
    add_footer(s, 3)

    # ============ 页4 Demo vs 产品 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "01 为什么需要 Harness", "从 Demo 到产品的鸿沟",
              "能跑的 Demo ≠ 可靠的产品")
    add_table(s, 0.8, 2.4, 11.7, [
        ["", "能跑的 Demo", "可靠的产品"],
        ["工具调用", "碰巧成功一次", "每次都正确编排"],
        ["异常处理", "出错即终止", "自动纠错恢复"],
        ["验证", "人肉看输出", "自动化断言"],
        ["关键", "LLM 能力展示", "Harness 工程保障"],
    ])
    add_textbox(s, 0.8, 5.4, 11.7, 0.8,
                "脆弱点正是 Harness 工程要解决的问题：约束防止越界、验证发现错误、纠正恢复异常",
                size=15, color=ACCENT, bold=True)
    add_footer(s, 4)

    # ============ 页5 核心公式 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "02 Harness 是什么", "核心公式",
              "最小公式是 Demo 视角，扩展公式是生产视角")
    add_textbox(s, 1.0, 2.3, 11, 0.8, "Agent = LLM + [ 上下文 + 工具 ]", size=30,
                bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(s, 1.0, 3.1, 11, 0.6, "扩展为生产形态：", size=16, color=GRAY,
                align=PP_ALIGN.CENTER)
    add_textbox(s, 1.0, 3.6, 11, 1.0,
                "Agent = LLM + [上下文 + 工具 + 约束 + 验证 + 纠正] = Model + Harness",
                size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_textbox(s, 1.0, 5.0, 11, 0.6,
                "把 LLM 当作核心组件 Model，围绕它构建的支撑代码统称 Harness",
                size=15, color=DARK, align=PP_ALIGN.CENTER)
    add_footer(s, 5)

    # ============ 页6 五要素拆解表 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "02 Harness 是什么", "五要素保障层", "每要素对应一个例子")
    add_table(s, 0.8, 2.3, 11.7, [
        ["要素", "作用", "例子（退款场景）"],
        ["上下文", "提供决策所需信息", "系统提示词写明 7 天退款政策"],
        ["工具", "提供执行能力", "query_order / process_refund"],
        ["约束", "限定能做什么", "校验退款金额 ≤ 订单金额"],
        ["验证", "检查做得对不对", "查数据库确认退款成功"],
        ["纠正", "做错了怎么办", "API 超时自动重试"],
    ])
    add_textbox(s, 0.8, 5.8, 11.7, 0.6,
                "同一模型，有无 Harness，结果天壤之别", size=15,
                color=ACCENT, bold=True)
    add_footer(s, 6)

    # ============ 页7 马具隐喻 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "02 Harness 是什么", "马具隐喻",
              "没有 Harness 的模型就像脱缰的野马")
    add_textbox(s, 1.0, 2.4, 11, 0.5, "野马（无 Harness）", size=20, bold=True,
                color=ACCENT)
    add_textbox(s, 1.0, 2.9, 11, 0.6,
                "能力惊人，但无法可靠地完成任务——不知道边界、不可预测", size=15, color=DARK)
    add_textbox(s, 1.0, 3.9, 11, 0.5, "装配马具（有 Harness）", size=20, bold=True,
                color=GREEN)
    add_textbox(s, 1.0, 4.4, 11, 0.6,
                "能力被引导、约束、放大为可靠的任务执行——可预测、可验证", size=15, color=DARK)
    add_textbox(s, 1.0, 5.4, 11, 0.6,
                "设计和优化这套模型之外基础设施的工程实践，就是 Harness 工程",
                size=16, bold=True, color=NAVY)
    add_footer(s, 7)

    # ============ 页8 LLM 编排工具 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "02 Harness 是什么", "LLM 编排工具",
              "模型负责选，框架负责跑、拦、纠")
    add_textbox(s, 1.0, 2.3, 11, 0.5, "工具（Tool）", size=18, bold=True, color=NAVY)
    add_textbox(s, 1.0, 2.8, 11, 0.5, "一个具体能力：bash、read、write、评估命令", size=14, color=DARK)
    add_textbox(s, 1.0, 3.4, 11, 0.5, "编排（Orchestration）", size=18, bold=True, color=NAVY)
    add_textbox(s, 1.0, 3.9, 11, 0.5,
                "如何让 LLM 正确使用工具：选哪个、顺序、参数、约束、纠错", size=14, color=DARK)
    add_textbox(s, 1.0, 4.6, 11, 0.5, "ReAct 循环", size=18, bold=True, color=ACCENT)
    add_textbox(s, 1.0, 5.1, 11, 1.0,
                "LLM 思考(Reason) → 选工具(Act) → 框架执行 → 结果回传 → 再思考\n"
                "Harness 在这一环介入：暴露哪些工具、允许/拒绝、校验参数、错误注入",
                size=14, color=DARK, line_spacing=1.2)
    add_footer(s, 8)

    # ============ 页9 验证独立性 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "02 Harness 是什么", "验证独立性",
              "被验证者不能给自己的试卷打分")
    add_textbox(s, 1.0, 2.4, 11, 0.5, "为什么断言不能交给 LLM？", size=20, bold=True,
                color=NAVY)
    add_textbox(s, 1.0, 3.0, 11, 1.5,
                "1. LLM 看不到完整事件流（工具调用序列在框架层）\n"
                "2. LLM 只能自我报告（'我执行了'），而 Harness 存在的意义就是不信自我报告\n"
                "3. 让被验证者调用验证者 = 自指循环，验证失去意义",
                size=15, color=DARK, line_spacing=1.3)
    add_textbox(s, 1.0, 5.2, 11, 0.8,
                "正确做法：断言由看得到事件流的一方执行 —— 框架插件 hook 或外部 runner",
                size=16, bold=True, color=ACCENT)
    add_footer(s, 9)

    # ============ 页10 概念小结 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "02 Harness 是什么", "概念小结",
              "两个视角，同一系统")
    add_table(s, 0.8, 2.4, 11.7, [
        ["", "最小公式（Demo）", "扩展公式（生产）"],
        ["组成", "LLM + 上下文 + 工具", "+ 约束 + 验证 + 纠正"],
        ["目的", "能跑起来", "长期可靠运转"],
        ["本质", "能力展示", "安全网兜底"],
    ])
    add_textbox(s, 0.8, 5.2, 11.7, 0.6,
                "扩展公式完全包含最小公式，并在外围加了一圈安全网", size=15,
                color=NAVY, bold=True)
    add_footer(s, 10)

    # ============ 页11 opencode 能力全景 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "03 怎么搭 Harness", "opencode 能力全景",
              "约束强、上下文强、验证弱 → 验证需自建")
    add_table(s, 0.8, 2.4, 11.7, [
        ["要素", "opencode 支持", "成熟度"],
        ["上下文", "Rules / Skills / References / Compaction hooks", "★★★★★"],
        ["工具", "内置14 + Custom Tools + MCP", "★★★★★"],
        ["约束", "Permissions / Policies / doom_loop", "★★★★★"],
        ["纠正", "Plugin hooks / revert / abort", "★★★★☆"],
        ["验证", "Structured Output（部分）", "★★☆☆☆"],
    ])
    add_textbox(s, 0.8, 5.8, 11.7, 0.5,
                "结论：验证层是最大短板，需要自建断言", size=15,
                color=ACCENT, bold=True)
    add_footer(s, 11)

    # ============ 页12 三层验证架构 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "03 怎么搭 Harness", "三层验证架构（本项目）",
              "单一事实源 + 分层断言")
    add_textbox(s, 1.0, 2.3, 11, 0.5, "1. 插件 hook（框架自动）", size=18, bold=True,
                color=NAVY)
    add_textbox(s, 1.0, 2.8, 11, 0.4,
                ".opencode/plugins/agent-assert.js → tool.execute.after 自动断言", size=13,
                color=DARK, mono=True)
    add_textbox(s, 1.0, 3.4, 11, 0.5, "2. 外部 runner（手动/CI）", size=18, bold=True,
                color=NAVY)
    add_textbox(s, 1.0, 3.9, 11, 0.4,
                "opencode-harness/harness_runner.py → 工具序列/越权/doom_loop", size=13,
                color=DARK, mono=True)
    add_textbox(s, 1.0, 4.5, 11, 0.5, "3. 规则单一事实源", size=18, bold=True,
                color=NAVY)
    add_textbox(s, 1.0, 5.0, 11, 0.4,
                "opencode-harness/rules.json → Python 与 JS 共享，避免漂移", size=13,
                color=DARK, mono=True)
    add_footer(s, 12)

    # ============ 页13 约束层落地 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "03 怎么搭 Harness", "约束层落地",
              "bash 白名单 + 只读评估 + 反死循环")
    add_code(s, 1.0, 2.3, 11.7, 3.2, """// opencode.json — architecture-quality agent 权限
"bash": {
  "*": "ask",                                    // 默认询问
  "arch-quality*": "allow",                      // 评估命令放行
  "python -m arch_quality*": "allow",
  "Test-Path*": "allow",                         // 只读探索
  "git status*": "allow", "git log*": "allow"
},
"edit": "deny",                                  // 评估任务只读
"write": "deny",
// doom_loop: 相同工具调用重复3次自动兜底""", size=12)
    add_footer(s, 13)

    # ============ 页14 验证与纠正 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "03 怎么搭 Harness", "验证与纠正",
              "断言器 + 重试，形成闭环")
    add_code(s, 1.0, 2.3, 11.7, 2.6, """// agent-assert.js — 评估命令后自动断言
"tool.execute.after": async (input, output) => {
  if (input.tool === "bash" && pendingEval) {
    const errors = runAsserts(output.output, rules);
    // 断言1: overall_score ∈ [0,100]
    // 断言2: 含 solver_physics 维度 或 明确非多物理场
    // 断言3: 无幻觉字段
    if (errors.length && rules.assert_mode === "throw")
      throw new Error(errors.join("; "));   // 纠正：LLM 收到错误后重试
  }
}""", size=12)
    add_textbox(s, 1.0, 5.2, 11.7, 0.6,
                "外部 runner 重试：超时或断言失败 → 新会话重跑（最多 N 次）", size=14,
                color=DARK)
    add_footer(s, 14)

    # ============ 页15 验证结果 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "04 效果与反思", "验证结果",
              "3 用例端到端 PASS + 发现真实 bug")
    add_table(s, 0.8, 2.3, 11.7, [
        ["用例", "目标", "结果"],
        ["case-1-multiphysics", "多物理场项目识别 + 评分", "PASS"],
        ["case-2-single-solver", "单求解器判定非多物理场", "PASS"],
        ["case-3-pure-markdown", "纯 Markdown 判定非多物理场", "PASS"],
    ])
    add_textbox(s, 0.8, 4.5, 11.7, 0.8,
                "额外收获：harness 在 case-2 中发现 arch_report.py 缺 ReportGenerator 导入的真实 bug",
                size=15, color=ACCENT, bold=True)
    add_textbox(s, 0.8, 5.3, 11.7, 0.5,
                "pytest 验证工具算得对，harness 验证 agent 用得好——两者互补", size=14,
                color=DARK)
    add_footer(s, 15)

    # ============ 页16 边界与未来 ============
    s = prs.slides.add_slide(blank)
    add_title(s, "04 效果与反思", "边界与未来",
              "验证层仍需自建，规则需单一事实源")
    add_textbox(s, 1.0, 2.3, 11, 0.5, "踩过的坑", size=18, bold=True, color=NAVY)
    add_textbox(s, 1.0, 2.8, 11, 1.2,
                "· SyntaxWarning 污染：Python 3.12+ 未加 r 前缀正则的警告刷屏\n"
                "· 非多物理场判据：单语言项目 JSON 无 solver_physics 提及\n"
                "· LLM 随机性：同 prompt 每次行为不同，需重试兜底",
                size=14, color=DARK, line_spacing=1.3)
    add_textbox(s, 1.0, 4.4, 11, 0.5, "未来方向", size=18, bold=True, color=NAVY)
    add_textbox(s, 1.0, 4.9, 11, 1.2,
                "· 接入 CI/CD 门禁（agent 变更后自动跑 harness）\n"
                "· 扩展真实项目用例（OpenFOAM 等）\n"
                "· 规则单一事实源继续演进（rules.json）",
                size=14, color=DARK, line_spacing=1.3)
    add_footer(s, 16)

    # 保存
    os.makedirs(os.path.dirname(OUT_PATH), exist_ok=True)
    prs.save(OUT_PATH)
    print(f"已生成: {OUT_PATH}")


if __name__ == "__main__":
    build()
